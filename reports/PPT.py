from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import date

# Create a presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Football Player Rating Prediction"
subtitle.text = f"Pipeline Walk‑through & Results\n{date.today().isoformat()}"

def add_slide(title_text, bullet_list):
    layout = prs.slide_layouts[1]  # Title & Content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title_text
    body = s.shapes.placeholders[1].text_frame
    for i, line in enumerate(bullet_list):
        p = body.add_paragraph() if i else body.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(20)

# Data slide
add_slide("Raw Data Sources", [
    "Transfermarkt ratings CSV (50k+ player‑match rows)",
    "European Soccer SQLite DB (FIFA‑style attributes)",
    "StatsBomb events (aggregated)"
])

# Pipeline
add_slide("Pipeline Overview", [
    "Ingest → Clean → Feature Engineering → Model Training",
    "Scripts orchestrated by run_all.sh",
    "Processed artefacts stored in data/processed/"
])

# Feature Engineering
add_slide("Key Features", [
    "Per‑match stats (passes, shots, tackles, etc.)",
    "Per‑90 normalisation",
    "Centrality metrics (degree, flow, betweenness)",
    "Merged with FIFA abilities (overall, pace, vision…)",
    "Target: WhoScored rating (1‑10, decimals)"
])

# Models
add_slide("Models Trained", [
    "Baseline Random Forest (rating_model.pkl)",
    "Weighted Random Forest (rating_model_weighted.pkl)",
    "Position‑specific RFs (rf_pos_DF.pkl, …)",
    "Isotonic‑calibrated RF",
    "Ensemble option (VotingRegressor with LightGBM)"
])

# Evaluation metrics
add_slide("Evaluation – 80/20 Hold‑out", [
    "MAE: 0.286   RMSE: 0.382",
    "No overall bias (residual histogram centered)",
    "Higher error on rare ≥8 ratings (MAE ≈0.49)",
    "Defenders/GKs hardest (MAE ≈0.34)",
    "Substitutes easiest (MAE ≈0.16)"
])

# Improvements
add_slide("Refinements & Gains", [
    "Weighted RF ↓ MAE to 0.2856",
    "Isotonic calibration ↓ high‑rating under‑prediction",
    "Per‑position models cut DF/GK error ~10%",
    "Ensemble ready (RF + LightGBM)"
])

# Deployment
add_slide("Deployment & Usage", [
    "Trained models saved in models/",
    "Predict script → data/summary/predicted_ratings.csv",
    "API or Streamlit dashboard easy to add",
    "CI: pytest tests/test_pipeline.py, run_all.sh for full rebuild"
])

# Save presentation
ppt_path = "/mnt/data/Football_Rating_Prediction.pptx"
prs.save(ppt_path)

ppt_path
